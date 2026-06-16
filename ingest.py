"""
ingest.py - 轻量版（纯SQLite存储，无需Embedding API）
上传时直接存文本，查询时用关键词匹配，避免API调用失败导致上传报错。
"""

import os
import io
import sqlite3
from pypdf import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

DB_PATH = os.environ.get("CHROMA_PATH", "/data/chroma_db").rstrip("/") + "/kb.sqlite3"


def _get_conn():
    os.makedirs(os.path.dirname(DB_PATH), exist_ok=True)
    conn = sqlite3.connect(DB_PATH)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS chunks (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            source TEXT NOT NULL,
            text TEXT NOT NULL
        )
    """)
    conn.commit()
    return conn


def extract_text_from_pdf(file_bytes: bytes) -> str:
    reader = PdfReader(io.BytesIO(file_bytes))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def extract_text_from_docx(file_bytes: bytes) -> str:
    doc = Document(io.BytesIO(file_bytes))
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    for table in doc.tables:
        for row in table.rows:
            text += "\n" + " | ".join(cell.text for cell in row.cells)
    return text


def extract_text_from_xlsx(file_bytes: bytes) -> str:
    wb = load_workbook(io.BytesIO(file_bytes), data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"【工作表: {sheet.title}】")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        slide_texts.append(line.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_texts.append(" | ".join(cells))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                slide_texts.append(f"[备注] {notes.strip()}")
        if slide_texts:
            parts.append(f"【幻灯片 {i}】\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


def extract_text_from_txt(file_bytes: bytes) -> str:
    for enc in ("utf-8", "gbk", "gb2312"):
        try:
            return file_bytes.decode(enc)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="ignore")


def extract_text_any(filename: str, file_bytes: bytes) -> str:
    name = filename.lower()
    if name.endswith(".pdf"):
        return extract_text_from_pdf(file_bytes)
    elif name.endswith(".docx"):
        return extract_text_from_docx(file_bytes)
    elif name.endswith((".xlsx", ".xlsm")):
        return extract_text_from_xlsx(file_bytes)
    elif name.endswith(".pptx"):
        return extract_text_from_pptx(file_bytes)
    elif name.endswith(".txt"):
        return extract_text_from_txt(file_bytes)
    else:
        raise ValueError("仅支持 .pdf、.docx、.xlsx、.pptx、.txt 文件")


def chunk_text(text: str, chunk_size: int = 800, overlap: int = 100) -> list:
    text = text.strip()
    if not text:
        return []
    chunks = []
    start = 0
    while start < len(text):
        chunk = text[start:start + chunk_size]
        if chunk.strip():
            chunks.append(chunk.strip())
        start += chunk_size - overlap
    return chunks


def ingest_document(filename: str, file_bytes: bytes, metadata: dict = None) -> dict:
    text = extract_text_any(filename, file_bytes)
    chunks = chunk_text(text)
    if not chunks:
        return {"filename": filename, "chunks_added": 0, "message": "未提取到文本内容"}

    conn = _get_conn()
    conn.execute("DELETE FROM chunks WHERE source = ?", (filename,))
    conn.executemany(
        "INSERT INTO chunks (source, text) VALUES (?, ?)",
        [(filename, chunk) for chunk in chunks]
    )
    conn.commit()
    conn.close()

    return {"filename": filename, "chunks_added": len(chunks)}


def query_knowledge_base(query: str, top_k: int = 5) -> list:
    conn = _get_conn()
    rows = conn.execute("SELECT source, text FROM chunks").fetchall()
    conn.close()

    if not rows:
        return []

    keywords = [w for w in query.replace("，", " ").replace("。", " ").split() if len(w) > 1]
    if not keywords:
        return [{"text": r[1], "source": r[0]} for r in rows[:top_k]]

    scored = []
    for source, text in rows:
        score = sum(text.count(kw) for kw in keywords)
        if score > 0:
            scored.append((score, source, text))

    scored.sort(key=lambda x: x[0], reverse=True)
    return [{"text": t, "source": s} for _, s, t in scored[:top_k]]


def list_documents() -> list:
    conn = _get_conn()
    rows = conn.execute("SELECT DISTINCT source FROM chunks").fetchall()
    conn.close()
    return [r[0] for r in rows]


def delete_document(filename: str) -> dict:
    conn = _get_conn()
    cursor = conn.execute("DELETE FROM chunks WHERE source = ?", (filename,))
    deleted = cursor.rowcount
    conn.commit()
    conn.close()
    return {"filename": filename, "deleted_chunks": deleted}
